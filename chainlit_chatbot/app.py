import os
import chainlit as cl
from openai import OpenAI
from dotenv import load_dotenv
from langchain.document_loaders import PyPDFLoader
from pptx import Presentation
import csv
import io
import logging
from typing import List, Dict, Tuple
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Set up logging for debugging purposes
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

# Load environment variables from .env file
load_dotenv()

# Initialize OpenAI client with API key from environment variables
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

@cl.on_chat_start
async def start():
    # Initialize an empty conversation history when a new chat starts
    cl.user_session.set("conversation_history", [])
    
    # Display a welcome message
    await cl.Message(content="Welcome to your AI assistant! How can I help you today?").send()


@cl.on_message
async def main(message: cl.Message):
    conversation_history: List[Dict[str, str]] = cl.user_session.get("conversation_history", [])
    
    async def process_file(file: cl.File) -> Tuple[str, str]:
        logger.info(f"Processing file: {file.name}")
        try:
            if file.name.lower().endswith('.pdf'):
                loader = PyPDFLoader(file.path)
                pages = loader.load_and_split()
                file_content = "\n\n".join(page.page_content for page in pages)
                logger.info(f"Extracted {len(pages)} pages from PDF, total length: {len(file_content)} characters")
                
                # Split the content into chunks
                text_splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=200)
                chunks = text_splitter.split_text(file_content)
                
                # Summarize each chunk
                chunk_summaries = []
                for i, chunk in enumerate(chunks):
                    chunk_summary = generate_summary(chunk)
                    chunk_summaries.append(f"Chunk {i+1} Summary: {chunk_summary}")
                
                # Combine chunk summaries
                combined_summary = "\n\n".join(chunk_summaries)
                
                # Generate a final summary of the combined summaries
                final_summary = generate_summary(combined_summary)
                
                logger.info(f"Generated summary for large PDF, final summary length: {len(final_summary)} characters")
            elif file.name.lower().endswith(('.ppt', '.pptx')):
                prs = Presentation(file.path)
                file_content = "\n\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
                logger.info(f"Extracted content from PPT, total length: {len(file_content)} characters")
            elif file.name.lower().endswith('.csv'):
                csv_content = []
                encodings = ['utf-8', 'iso-8859-1', 'windows-1252']
                for encoding in encodings:
                    try:
                        with open(file.path, 'r', newline='', encoding=encoding) as csvfile:
                            csv_reader = csv.reader(csvfile)
                            for row in csv_reader:
                                csv_content.append(", ".join(row))
                        file_content = "\n".join(csv_content)
                        logger.info(f"Extracted content from CSV using {encoding} encoding, total length: {len(file_content)} characters")
                        break
                    except UnicodeDecodeError:
                        if encoding == encodings[-1]:
                            raise ValueError(f"Unable to decode CSV file with any of the attempted encodings: {', '.join(encodings)}")
                        continue
            else:
                raise ValueError("Unsupported file type. Please upload a PDF, PPT, or CSV file.")

            if not file_content.strip():
                raise ValueError("No content could be extracted from the file.")
        
            return file_content
        except Exception as e:
            logger.error(f"Error processing file: {str(e)}", exc_info=True)
            raise

    if message.elements:
        for element in message.elements:
            if isinstance(element, cl.File) and (element.mime in ["application/pdf", "application/vnd.openxmlformats-officedocument.presentationml.presentation", "text/csv"] or element.name.lower().endswith('.csv')):
                try:
                    file_content, summary = await process_file(element)
                    file_type = "PDF" if element.mime == "application/pdf" else "PPT" if element.mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation" else "CSV"
                    conversation_history.append({"role": "system", "content": f"{file_type} Content: {file_content}\n\nSummary: {summary}"})
                    await cl.Message(content=f"📄 File '{element.name}' processed. Here's a summary:\n\n{summary}\n\nYou can now ask questions about this document.").send()
                except Exception as e:
                    await cl.Message(content=f"❌ Error processing file: {str(e)}").send()
                return
            else:
                await cl.Message(content="❌ Unsupported file type. Please upload a PDF, PPT, or CSV file.").send()
                return

    conversation_history.append({"role": "user", "content": message.content})

    response = client.chat.completions.create(
        model="gpt-4",
        messages=conversation_history,
        max_tokens=300
    )

    assistant_message = response.choices[0].message
    conversation_history.append({"role": "assistant", "content": assistant_message.content})
    cl.user_session.set("conversation_history", conversation_history)

    await cl.Message(content=assistant_message.content).send()

def generate_summary(text):
    """
    Generate a summary of the given text using OpenAI's GPT-4 model.
    
    :param text: The text to summarize
    :return: A summary of the text
    """
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that summarizes documents."},
            {"role": "user", "content": f"Please summarize the following text in about 3-4 sentences:\n\n{text}"}
        ],
        max_tokens=150
    )
    return response.choices[0].message.content

if __name__ == "__main__":
    # Run the Chainlit app
    cl.run()
